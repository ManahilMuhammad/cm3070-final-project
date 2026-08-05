import io
import json
import pytest
import instrumentation as inst
import pipeline

@pytest.fixture(autouse=True)
def clean_records():
    inst.reset()
    yield
    inst.reset()


def chat_response(content):
    return {"message": {"content": content}}


# --> fuse()

def test_fuse_combines_sections():
    result = pipeline.fuse("t", "s", "n", "f")
    assert result == "[TRANSCRIPT]\nt\n\n[SLIDE]\ns\n\n[NOTES]\nn\n\n[FIGURE]\nf"


def test_fuse_skips_empty():
    result = pipeline.fuse("t", "", "  ", "None")
    assert result == "[TRANSCRIPT]\nt"


def test_fuse_is_timed_by_instrumentation():
    pipeline.fuse("t", "", "", "")
    assert "fusion of extracted text" in inst.aggregate()


# --> extract_slides()

def test_extract_slides_none():
    assert pipeline.extract_slides(None) == ""


def test_extract_slides_unsupported():
    fake = io.BytesIO(b"data")
    fake.name = "notes.txt"
    assert pipeline.extract_slides(fake) == ""


def test_extract_slides_reads_pdf_text(tmp_path):
    fitz = pytest.importorskip("fitz")
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "hello world from pdf")
    data = doc.tobytes()
    doc.close()

    fake = io.BytesIO(data)
    fake.name = "slides.pdf"

    text = pipeline.extract_slides(fake)
    assert "hello world from pdf" in text


def test_extract_slides_reads_pptx_text():
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "hello world from pptx"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    buf.name = "slides.pptx"

    text = pipeline.extract_slides(buf)
    assert "hello world from pptx" in text


# --> describe_figure()

def test_describe_figure_none():
    assert pipeline.describe_figure(None) == ""


def test_describe_figure_calls_ollama_chat(monkeypatch):
    captured = {}

    def fake_chat(model, messages):
        captured["model"] = model
        captured["images"] = messages[0]["images"]
        return chat_response("a diagram of X")

    monkeypatch.setattr(pipeline.ollama, "chat", fake_chat)

    fake = io.BytesIO(b"imgbytes")
    result = pipeline.describe_figure(fake)

    assert result == "a diagram of X"
    assert captured["model"] == "qwen2.5vl:latest"
    assert captured["images"] == [b"imgbytes"]
    assert "figure description" in inst.aggregate()


# --> make_summary()

def test_make_summary_calls_ollama_chat(monkeypatch):
    monkeypatch.setattr(
        pipeline.ollama, "chat", lambda **kw: chat_response("# Summary\n...")
    )

    result = pipeline.make_summary("some combined text")

    assert result == "# Summary\n..."
    assert "summarisation" in inst.aggregate()


def test_make_summary_deterministic_opts(monkeypatch):
    captured = {}

    def fake_chat(model, messages, options=None):
        captured["messages"] = messages
        captured["options"] = options
        return chat_response("# Summary\n...")

    monkeypatch.setattr(pipeline.ollama, "chat", fake_chat)

    pipeline.make_summary("some combined text")

    assert captured["options"] == {"temperature": 0.2, "seed": 42}
    assert captured["messages"][0]["role"] == "system"
    assert captured["messages"][1]["role"] == "user"
    assert "some combined text" in captured["messages"][1]["content"]


# --> generate_ques() / make_quiz()

def test_generate_ques_succeeds_first_try(monkeypatch):
    valid = json.dumps({"questions": [{"question": "2+2?", "answer": "4", "options": ["3", "4", "5", "6"]}]})
    monkeypatch.setattr(pipeline.ollama, "chat", lambda **kw: chat_response(valid))

    spec = {"type": "multiple-choice", "shape": {}, "instructions": ""}
    qs = pipeline.generate_ques("content", spec, ["Topic A"])

    assert len(qs) == 1
    assert qs[0]["type"] == "multiple-choice"
    assert qs[0]["question"] == "2+2?"
    assert qs[0]["answer"] == "4"
    assert qs[0]["topic"] == "Topic A"


def test_generate_ques_retries_only_failed_items(monkeypatch):
    calls = []

    def fake_chat(**kw):
        calls.append(kw)
        if len(calls) == 1:
            # 2 requested, first valid, second invalid (missing answer)
            return chat_response(json.dumps({"questions": [
                {"question": "Q about A", "answer": "a1", "options": None},
                {"question": "bad", "answer": None, "options": None},
            ]}))
        # second call should only be asked for the one that failed
        return chat_response(json.dumps({"questions": [
            {"question": "Q about B", "answer": "a2", "options": None},
        ]}))

    monkeypatch.setattr(pipeline.ollama, "chat", fake_chat)

    spec = {"type": "short-answer", "shape": {}, "instructions": ""}
    qs = pipeline.generate_ques("content", spec, ["Topic A", "Topic B"], retries=3)

    assert len(calls) == 2
    assert '"Topic B"' in calls[1]["messages"][1]["content"]
    assert '"Topic A"' not in calls[1]["messages"][1]["content"]
    assert [q["question"] for q in qs] == ["Q about A", "Q about B"]
    assert qs[1]["topic"] == "Topic B"


def test_generate_ques_blank_answer_leak(monkeypatch):
    # answer word appears in the question -> invalid -> exhausts retries -> dropped
    bad = json.dumps({"questions": [{"question": "Paris is the _____ of France", "answer": "Paris", "options": None}]})
    monkeypatch.setattr(pipeline.ollama, "chat", lambda **kw: chat_response(bad))

    spec = {"type": "fill-in-the-blank", "shape": {}, "instructions": ""}
    qs = pipeline.generate_ques("content", spec, [None], retries=2)

    assert qs == []


def test_true_false_batch_rejects_question_phrasing(monkeypatch):
    # phrased as question instead of statement -> invalid -> dropped
    bad = json.dumps({"questions": [{"question": "Which process does the brain use to learn?", "answer": "true", "options": None}]})
    monkeypatch.setattr(pipeline.ollama, "chat", lambda **kw: chat_response(bad))

    spec = {"type": "true-false", "shape": {}, "instructions": ""}
    qs = pipeline.generate_ques("content", spec, [None], retries=2)

    assert qs == []


def test_true_false_batch_accepts_statement(monkeypatch):
    good = json.dumps({"questions": [{"question": "The brain uses backpropagation to learn.", "answer": "false", "options": None}]})
    monkeypatch.setattr(pipeline.ollama, "chat", lambda **kw: chat_response(good))

    spec = {"type": "true-false", "shape": {}, "instructions": ""}
    qs = pipeline.generate_ques("content", spec, [None], retries=2)

    assert len(qs) == 1
    assert qs[0]["answer"] == "false"


def test_make_quiz_filters_failed(monkeypatch):
    monkeypatch.setattr(pipeline, "extract_topics", lambda *a, **k: ["Topic A", "Topic B"])
    calls = {"n": 0}

    def fake_batch(content, spec, topics, retries=4, avoid=None):
        calls["n"] += 1
        if spec["type"] == "true-false":
            return [] # both failed and got dropped
        return [{"type": spec["type"], "question": f"q{calls['n']}-{i}", "answer": "a", "options": None, "topic": t} for i, t in enumerate(topics)]

    monkeypatch.setattr(pipeline, "generate_ques", fake_batch)

    quiz = pipeline.make_quiz("combined text")

    assert len(quiz) == 6
    assert all(q["type"] != "true-false" for q in quiz)
    assert "create quiz" in inst.aggregate()


def test_make_quiz_two_per_type(monkeypatch):
    monkeypatch.setattr(pipeline, "extract_topics", lambda *a, **k: ["Topic A", "Topic B"])

    def fake_batch(content, spec, topics, retries=4, avoid=None):
        return [{"type": spec["type"], "question": f"{spec['type']}-{i}", "answer": "a", "options": None, "topic": t} for i, t in enumerate(topics)]

    monkeypatch.setattr(pipeline, "generate_ques", fake_batch)

    quiz = pipeline.make_quiz("combined text")

    assert len(quiz) == 8
    counts = {}
    for q in quiz:
        counts[q["type"]] = counts.get(q["type"], 0) + 1
    assert all(c == 2 for c in counts.values())


def test_make_quiz_passes_avoid_list(monkeypatch):
    monkeypatch.setattr(pipeline, "extract_topics", lambda *a, **k: ["Topic A", "Topic B"])
    seen_avoid = []

    def fake_batch(content, spec, topics, retries=4, avoid=None):
        seen_avoid.append(list(avoid or []))
        return [{"type": spec["type"], "question": f"{spec['type']}-q{i}", "answer": "a", "options": None, "topic": t} for i, t in enumerate(topics)]

    monkeypatch.setattr(pipeline, "generate_ques", fake_batch)

    pipeline.make_quiz("combined text")

    # first type's batch has nothing to avoid
    # second type's batch is told about every question from the first batch
    assert seen_avoid[0] == []
    assert seen_avoid[1] == ["multiple-choice-q0", "multiple-choice-q1"]


def test_make_quiz_topics_round_robin(monkeypatch):
    monkeypatch.setattr(pipeline, "extract_topics", lambda *a, **k: ["Topic A", "Topic B"])
    seen_topics = []

    def fake_batch(content, spec, topics, retries=4, avoid=None):
        seen_topics.extend(topics)
        return [{"type": spec["type"], "question": f"q-{len(seen_topics)}-{i}", "answer": "a", "options": None, "topic": t} for i, t in enumerate(topics)]

    monkeypatch.setattr(pipeline, "generate_ques", fake_batch)

    quiz = pipeline.make_quiz("combined text")

    assert seen_topics == ["Topic A", "Topic B"] * 4
    assert [q["topic"] for q in quiz] == seen_topics


def test_make_quiz_no_topics_fallback(monkeypatch):
    monkeypatch.setattr(pipeline, "extract_topics", lambda *a, **k: [])
    monkeypatch.setattr(
        pipeline, "generate_ques",
        lambda content, spec, topics, retries=4, avoid=None: [
            {"type": spec["type"], "question": f"q-{i}", "answer": "a", "options": None, "topic": t or "General"}
            for i, t in enumerate(topics)
        ]
    )

    quiz = pipeline.make_quiz("combined text")

    assert all(q["topic"] == "General" for q in quiz)


# --> extract_topics()

def test_extract_topics_dedupes(monkeypatch):
    raw = json.dumps({"topics": ["Neural Nets", "neural nets", "Backprop"]})
    monkeypatch.setattr(pipeline.ollama, "generate", lambda **kw: {"response": raw})

    topics = pipeline.extract_topics("lecture text", n=8)

    assert topics == ["Neural Nets", "Backprop"]


def test_extract_topics_bad_json_fallback(monkeypatch):
    monkeypatch.setattr(pipeline.ollama, "generate", lambda **kw: {"response": "not json"})

    topics = pipeline.extract_topics("lecture text")

    assert topics == ["General"]


# --> generate_score()

def test_generate_score_correctness():
    results = [
        {"time": 4.0, "num_words": 4, "user_answer": "Paris", "correct_answer": "paris", "topic": "geo"},
        {"time": 2.0, "num_words": 4, "user_answer": "wrong", "correct_answer": "right", "topic": "geo"},
    ]

    out = json.loads(pipeline.generate_score(results))

    assert results[0]["correct"] is True
    assert results[1]["correct"] is False
    assert results[1]["confidence"] == 0.0
    assert 0.0 <= results[0]["confidence"] <= 1.0

    assert len(out) == 1
    assert out[0]["topic"] == "geo"
    assert out[0]["score"] == "1/2"
    assert "generate score" in inst.aggregate()


def test_generate_score_zero_stdev():
    results = [{"time": 1.0, "num_words": 2, "user_answer": "a", "correct_answer": "a", "topic": "t"}]

    pipeline.generate_score(results)

    assert results[0]["correct"] is True
    assert results[0]["confidence"] == pytest.approx(0.5)


def test_generate_score_sorts_topics():
    results = [
        {"time": 1.0, "num_words": 2, "user_answer": "a", "correct_answer": "a", "topic": "strong"},
        {"time": 1.0, "num_words": 2, "user_answer": "x", "correct_answer": "y", "topic": "weak"},
    ]

    out = json.loads(pipeline.generate_score(results))

    assert [o["topic"] for o in out] == ["weak", "strong"]


# --> generate_feedback() / create_plan()

def test_generate_feedback_calls_ollama_chat(monkeypatch):
    monkeypatch.setattr(pipeline.ollama, "chat", lambda **kw: chat_response("Great job!"))

    result = pipeline.generate_feedback("[]", "combined")

    assert result == "Great job!"
    assert "generate feedback" in inst.aggregate()


def test_feedback_has_system_message(monkeypatch):
    captured = {}

    def fake_chat(model, messages):
        captured["messages"] = messages
        return chat_response("Great job!")

    monkeypatch.setattr(pipeline.ollama, "chat", fake_chat)

    pipeline.generate_feedback("[]", "combined")

    assert captured["messages"][0]["role"] == "system"
    assert captured["messages"][1]["role"] == "user"


def test_create_plan_empty_skips_llm(monkeypatch):
    called = []
    monkeypatch.setattr(pipeline.ollama, "chat", lambda **kw: called.append(1) or chat_response("{}"))

    plan = pipeline.create_plan("[]", "combined")

    assert plan == []
    assert called == []
    assert "create plan" in inst.aggregate()


def test_create_plan_more_days_for_weak(monkeypatch):
    performance = json.dumps([
        {"topic": "Weak Topic", "score": "0/2", "confidence": 0.1},
        {"topic": "Strong Topic", "score": "2/2", "confidence": 0.9},
    ])

    monkeypatch.setattr(
        pipeline.ollama, "chat",
        lambda **kw: chat_response(json.dumps({
            "actions": {
                "Weak Topic": ["Re-read section A", "Practice problem set B"],
                "Strong Topic": ["Quick review of key terms"],
            }
        }))
    )

    plan = pipeline.create_plan(performance, "combined", duration_days=6)

    assert [item["day"] for item in plan] == [1, 2, 3, 4, 5, 6]
    weak_days = [item for item in plan if item["topic"] == "Weak Topic"]
    strong_days = [item for item in plan if item["topic"] == "Strong Topic"]
    assert len(weak_days) > len(strong_days)
    assert weak_days[0]["priority"] == "high"
    assert strong_days[0]["priority"] == "low"


def test_create_plan_short_duration(monkeypatch):
    performance = json.dumps([
        {"topic": "A", "confidence": 0.1},
        {"topic": "B", "confidence": 0.3},
        {"topic": "C", "confidence": 0.9},
    ])
    monkeypatch.setattr(
        pipeline.ollama, "chat",
        lambda **kw: chat_response(json.dumps({"actions": {"A": ["a1"], "B": ["b1"], "C": ["c1"]}}))
    )

    plan = pipeline.create_plan(performance, "combined", duration_days=2)

    assert [item["topic"] for item in plan] == ["A", "B"]


def test_allocate_days_weights_confidence():
    topics = [("A", 0.1), ("B", 0.5), ("C", 0.9)]

    allocation = pipeline._allocate_days(topics, 7)

    assert sum(days for _, _, days in allocation) == 7
    days_by_topic = {name: days for name, _, days in allocation}
    assert days_by_topic["A"] > days_by_topic["B"] > days_by_topic["C"]


def test_allocate_days_even_split_weighted():
    # since duration is divided evenly across topics
    # a remainder-only split would ignore confidence
    # proportional weighting must not
    topics = [("Weak", 0.1), ("Strong", 0.9)]

    allocation = pipeline._allocate_days(topics, 6)

    days_by_topic = {name: days for name, _, days in allocation}
    assert sum(days_by_topic.values()) == 6
    assert days_by_topic["Weak"] > days_by_topic["Strong"]


def test_priority_bucket_thresholds():
    assert pipeline._priority_bucket(0.1) == "high"
    assert pipeline._priority_bucket(0.5) == "medium"
    assert pipeline._priority_bucket(0.9) == "low"


# --> unload_all()

def test_unload_all_clears_models():
    pipeline._models["dummy"] = object()

    pipeline.unload_all()

    assert pipeline._models == {}


# --> release_llm() / end_of_phase wiring into instrumentation

def test_release_llm_error_skips_inst(monkeypatch):
    def fake_generate(model, prompt, keep_alive):
        raise RuntimeError("ollama down")

    monkeypatch.setattr(pipeline.ollama, "generate", fake_generate)

    calls = []
    monkeypatch.setattr(pipeline.inst, "report", lambda *a, **k: calls.append("report"))
    monkeypatch.setattr(pipeline.inst, "save_run", lambda *a, **k: calls.append("save_run"))
    monkeypatch.setattr(pipeline.inst, "reset", lambda *a, **k: calls.append("reset"))

    ok = pipeline.release_llm()

    assert ok is False
    assert calls == []


def test_release_llm_end_of_phase_true(monkeypatch):
    monkeypatch.setattr(pipeline.ollama, "generate", lambda model, prompt, keep_alive: {})

    calls = []
    monkeypatch.setattr(pipeline.inst, "report", lambda *a, **k: calls.append(("report", a, k)))
    monkeypatch.setattr(pipeline.inst, "save_run", lambda *a, **k: calls.append(("save_run", a, k)))
    monkeypatch.setattr(pipeline.inst, "reset", lambda *a, **k: calls.append(("reset", a, k)))

    ok = pipeline.release_llm(label="my-run", end_of_phase=True)

    assert ok is True
    assert [c[0] for c in calls] == ["report", "save_run", "reset"]
    assert calls[1][1] == ("my-run",)


def test_release_llm_end_of_phase_false(monkeypatch):
    monkeypatch.setattr(pipeline.ollama, "generate", lambda model, prompt, keep_alive: {})

    calls = []
    monkeypatch.setattr(pipeline.inst, "report", lambda *a, **k: calls.append("report"))
    monkeypatch.setattr(pipeline.inst, "save_run", lambda *a, **k: calls.append("save_run"))
    monkeypatch.setattr(pipeline.inst, "reset", lambda *a, **k: calls.append("reset"))

    ok = pipeline.release_llm(model="qwen2.5vl:latest", end_of_phase=False)

    assert ok is True
    assert calls == []


def test_release_llm_integration(monkeypatch, tmp_path):
    monkeypatch.setattr(pipeline.ollama, "generate", lambda model, prompt, keep_alive: {})

    log_path = tmp_path / "timings.jsonl"
    real_save_run = inst.save_run
    monkeypatch.setattr(
        pipeline.inst, "save_run",
        lambda label, notes="", path=log_path: real_save_run(label, notes, path)
    )

    pipeline.fuse("transcript text", "", "", "")
    assert "fusion of extracted text" in inst.aggregate()

    ok = pipeline.release_llm(label="integration-run", end_of_phase=True)

    assert ok is True
    assert inst._records == []

    lines = log_path.read_text(encoding="utf-8").strip().splitlines()
    assert len(lines) == 1
    record = json.loads(lines[0])
    assert record["label"] == "integration-run"
    assert "fusion of extracted text" in record["stages"]

