import ollama
import json
from faster_whisper import WhisperModel
from transformers import TrOCRProcessor, VisionEncoderDecoderModel
from PIL import Image
import cv2
import numpy as np
import fitz
from pptx import Presentation
from collections import defaultdict
import io
import streamlit as st
import instrumentation as inst
import statistics, math, json
from collections import defaultdict
import gc
import torch

_models = {}

def _get(name, loader):
    """
    loading cache for Whisper, TrOCR
    """
    if name not in _models:
        _models[name] = loader()
    return _models[name]

def unload_all():
    """ 
    releases the ingestion models
    called once extraction is finished 
    """
    _models.clear()
    gc.collect()
    if torch.cuda.is_available():
        torch.cuda.empty_cache()

DEVICE = "cuda" if torch.cuda.is_available() else "cpu"

def _load_trocr():
    """
    loads TrOCR once
    """
    name = "microsoft/trocr-base-handwritten"
    processor = TrOCRProcessor.from_pretrained(name)
    model = VisionEncoderDecoderModel.from_pretrained(name).to(DEVICE)
    model.eval()
    return processor, model

@inst.timed("transcription")
def transcribe_audio(path):
    if path is None:
        return ""

    whisper_model = _get("whisper", lambda: WhisperModel("base", compute_type="int8"))
    
    segments, info = whisper_model.transcribe(path)
    transcript = " ".join(segment.text for segment in segments)

    return transcript

@inst.timed("slide extraction")
def extract_slides(path):
    if path is None:
        return ""
    
    name = path.name.lower()

    # PDF
    if name.endswith('.pdf'):
        data = path.read()
        doc = fitz.open(stream=data, filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text

    # PPTX
    elif name.endswith('.pptx'):
        prs = Presentation(path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = ''.join(run.text for run in para.runs)
                        if text.strip():
                            lines.append(text)
        return '\n'.join(lines)

    # other / incorrect file types do not get parsed here
    else:
        return ""

@inst.timed("OCR")
def ocr_notes(path):
    if path is None:
        return ""

    trocr_processor, trocr_model = _get("trocr", _load_trocr)

    # minimum size of one box
    min_area = 0.0004
    pad = 6

    image = Image.open(path).convert("RGB")
    img = np.array(image.convert("L")) # convert to grayscale
    h, w = img.shape # get height and width
    mean_brightness = img.mean() # get brightness

    # if background is dark with light text then do not invert
    if mean_brightness < 127:         
        thresh_type = cv2.THRESH_BINARY  

    # if background is light with dark text then invert     
    else:                            
        thresh_type = cv2.THRESH_BINARY_INV   

    _, binary = cv2.threshold(img, 0, 255, thresh_type + cv2.THRESH_OTSU)

    # dilate white pixels horizontally so letters in a single line merge
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (max(20, w // 20), 5))
    dilated = cv2.dilate(binary, kernel, iterations=1)

    # get shape outlines for the merged white pixels
    contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    boxes = [cv2.boundingRect(c) for c in contours] # create bounding boxes around the shapes
    boxes = [b for b in boxes if b[2] * b[3] > min_area * w * h] # only keep boxes larger than the minimum area
    boxes.sort(key=lambda b: b[1]) # sort boxes from top to bottom

    # crop each box into a separate image
    lines = []
    for x, y, bw, bh in boxes:
        lines.append(image.crop((max(0, x - pad), max(0, y - pad),
                                    min(w, x + bw + pad), min(h, y + bh + pad))))

    # apply TrOCR on each image
    notes = ""
    for line in lines:
        pixel_values = trocr_processor(images=line.convert("RGB"), return_tensors="pt").pixel_values
        generated_ids = trocr_model.generate(pixel_values)
        notes += trocr_processor.batch_decode(generated_ids, skip_special_tokens=True)[0] + "\n" # append extracted text

    return notes

@inst.timed("figure description")
def describe_figure(path): 
    prompt = """
You are helping create study notes from lecture material.
Describe the figure, chart, or diagram in the uploaded image clearly including all key data,
trends, labels, axes, or relationships, so it can be understood without seeing the image.
If the image contains no figure or diagram, reply with "None".
"""

    if path is None:
        return ""

    path.seek(0)
    image = path.read()

    description = ollama.chat(
        model="qwen2.5vl:latest",
        messages=[{"role": "user", "content": prompt, "images": [image]}],
    )['message']['content']

    return description

@inst.timed("fusion of extracted text")
def fuse(transcript, slides, notes, figure):
    """
    combines all pieces of extracted text together
    """

    parts = []
    components = [
        {'name': "TRANSCRIPT", 'content': transcript}, 
        {'name': "SLIDE", 'content': slides}, 
        {'name': "NOTES", 'content': notes}, 
        {'name': "FIGURE", 'content': figure}
    ]

    for component in components:
        content = component['content'].strip()
        if content and content.lower() != "none":
            parts.append(f"[{component['name']}]\n{content}")

    fusion = '\n\n'.join(parts)

    return fusion

@inst.timed("summarisation")
def make_summary(combined):
    system = (
        "You are a careful academic note-taker. You turn lecture material into "
        "structured study notes. You never invent facts, examples, or figures "
        "that are not present in the source text. You output ONLY the notes "
        "themselves - never a preamble, never a note to the person who asked, "
        "and never a sign-off."
    )

    prompt = f"""Summarise the lecture material below into structured study notes.

Follow this exact structure:
1. A single "# Overview" section: 2-3 sentences describing what the material covers.
2. One "## <Topic>" heading for EACH distinct topic, in the order it first appears.
3. Under each topic heading, 3-6 concise bullet points with the key facts, definitions, and relationships for that topic.
4. A final "## Key Terms" section: a bulleted list formatted as "**Term** - one-line definition".

Rules:
- Use ONLY information present in the TEXT below. Do not add outside knowledge or examples not present in the text.
- Mention every topic covered in the TEXT; do not omit or merge distinct topics.
- Be concise and factual. No filler like "In this lecture" or "This section discusses".
- Use consistent markdown: "#" for the title, "##" for headings, "-" for bullets.
- Do NOT include any preamble, meta-commentary, or sign-off - no "Here is a summary", no "Based on the text provided", no "I hope this helps". Output must start immediately with "# Overview" and end after the Key Terms section.

TEXT:
{combined}
"""

    summary = ollama.chat(
        model="llama3.2",
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
        options={"temperature": 0.2, "seed": 42},
    )['message']['content']

    return summary

# words that cannot be used as blank words in fill-in-the-blank
_BLANK_STOPWORDS = {
    "a", "an", "the", "is", "are", "was", "were", "be", "been", "being",
    "of", "in", "on", "at", "to", "for", "and", "or", "but", "it", "its",
    "this", "that", "these", "those", "as", "by", "with", "from", "not",
}

# words that cannot appear at the beginning of a true-false question 
# (because they would make it a question rather than a statement that is true or false)
_QUESTION_STARTERS = {
    "which", "what", "who", "whom", "whose", "how", "why", "when", "where",
    "does", "do", "did", "is", "are", "was", "were", "can", "could", "will",
    "would", "should", "has", "have", "had",
}

def _question_words(text):
    """
    strips the question to only important words therein
    """
    return {w.strip(".,?!:;'\"()_").lower() for w in text.split() if len(w.strip(".,?!:;'\"()_")) > 2}

def _too_similar(question, asked, threshold=0.5):
    """
    checks for duplicates
    high word overlap with any previously asked question
    """
    q_words = _question_words(question) # simplify question to root words

    if not q_words:
        return False

    # search previously asked questions and see if any of them 
    # majorly overlap with this one
    for prev in asked:
        prev_words = _question_words(prev)
        if not prev_words:
            continue
        overlap = len(q_words & prev_words) / len(q_words | prev_words)
        if overlap >= threshold:
            return True
        
    return False

def _valid_answer(ques_type, question, answer, options, avoid=None):
    """
    retries bad generation of each type of question instead of returning it
    """

    ans = str(answer).strip()

    # if detected to be a duplicate then invalid
    if _too_similar(question, avoid or []):
        return False

    # MCQ
    if ques_type == "multiple-choice":

        # check for correct type and length of options list
        if not isinstance(options, list) or len(options) != 4:
            return False
        
        cleaned = [str(o).strip() for o in options if str(o).strip()]

        # check that each MCQ option is distinct
        if len(cleaned) != 4 or len({o.lower() for o in cleaned}) != 4:
            return False
        
        return ans in cleaned

    # fill-in-the-blank
    if ques_type == "fill-in-the-blank":

        # check for missing blank
        if "_____" not in question:
            return False

        # check that answer is not present within question
        if ans.lower() in question.lower():
            return False

        # check that answer is a valid word
        return len(ans) >= 3 and ans.lower() not in _BLANK_STOPWORDS

    # true-false
    if ques_type == "true-false":

        # check that answer is either true or false
        if ans.lower() not in ("true", "false"):
            return False

        stripped = question.strip()

        # check that the sentence is not a question
        if stripped.endswith("?"):
            return False

        # check that the first word is not one of the question starters
        first_word = stripped.split(" ", 1)[0].lower().strip(",.:;\"'")
        return first_word not in _QUESTION_STARTERS

    # short-answer
    if ques_type == "short-answer":
        # check that answer is between 1 and 6 words, inclusive (short)
        return 1 <= len(ans.split()) <= 6

    return True

def generate_ques(content, spec, retries=4, avoid=None, topic=None):
    avoid = avoid or []
    avoid_block = ""
    if avoid:
        already = "\n".join(f"- {a}" for a in avoid)
        avoid_block = f"""
Do NOT repeat or closely rephrase any of these questions already asked. Ask about a DIFFERENT fact or aspect of the content:
{already}
""" # avoid duplicates

    # give the topic to make the question about
    topic_block = ""
    if topic:
        topic_block = f'\nFocus specifically on this topic: "{topic}". Do not ask about any other topic.\n'

    # PROMPT
    prompt = f"""You are a quiz generator. Read the content, then create ONE {spec["type"]} question.

CONTENT:
{content}

Now generate ONE {spec["type"]} question about the content above.
{spec["instructions"]}
{topic_block}{avoid_block}
General rules:
- Base the question ONLY on facts stated in the content. Do not invent facts, numbers, or terms that are not in the content.
- Write an original question in your own words. Do NOT copy a sentence verbatim from the content as the question.
- The question must be understandable and answerable on its own, without needing to see the original content.

Return ONLY a JSON object with these EXACT keys: "question", "answer", "options".
Example shape: {json.dumps(spec["shape"])}
Do NOT return a summary. Do NOT use any other keys."""

    ques_type = spec["type"]
    q = {}

    # retry 4 times to get valid question
    for retry in range(retries):

        # generate question
        response = ollama.chat(model="llama3.2",
                               messages=[
                                   {"role": "system", "content": "You generate quiz questions as JSON. You never summarise."},
                                   {"role": "user", "content": prompt}
                                ],
                                format="json",
                                options={"temperature": 0.3},
                                )

        try:
            q = json.loads(response["message"]["content"])
        except json.JSONDecodeError:
            print(f"Failed at attempt {retry}")
            continue

        ques = q.get("question")
        ans = q.get("answer")

        # validate questions and answers and return if valid
        if ques and ans not in (None, "", []) and _valid_answer(ques_type, ques, ans, q.get("options"), avoid):
            return {
                "type": ques_type,
                "question": ques,
                "answer": ans,
                "options": q.get("options")
                }

    # return a failed question if all retries exhausted
    return {
            "type": ques_type,
            "question": q.get("question"),
            "answer": None,
            "options": None,
            "failed": True
            }

@inst.timed("create quiz")
def make_quiz(combined, per_type=2, topics=None):
    questions = []

    # extract topics if not already done
    if topics is None:
        topics = extract_topics(combined, n=8)

    # specifications for each question type
    qspecs = [
        {
            "type": "multiple-choice",
            "shape": {"question": "<question>", "answer": "<correct answer>", "options": ["<option1>","<option2>","<option3>","<option4>"]},
            "instructions": (
                "Ask about a specific fact, definition, or relationship from the content. "
                "Provide exactly 4 options. Exactly ONE must be correct and explicitly "
                "supported by the content; the other 3 must be plausible but clearly "
                "wrong - not just a reworded version of the correct answer, and not "
                "also technically correct. All 4 options must be distinct from each "
                "other and similar in length and style. 'answer' MUST be copied "
                "exactly, character-for-character, from one of the 'options'."
            )
        },
        {
            "type": "fill-in-the-blank",
            "shape": {"question": "<sentence with a missing word represented as _____>", "answer": "<missing word>", "options": None},
            "instructions": (
                "Write ONE original factual sentence based on the content (do not "
                "copy a sentence verbatim from the content). Replace ONE important, "
                "specific keyword or technical term that is central to that "
                "sentence's meaning with '_____'. NEVER blank a common word like "
                "'the', 'is', 'a', or 'this' - the blanked word must be something a "
                "student who understood the topic could infer from the rest of the "
                "sentence. The 'question' is the sentence WITH the blank. The "
                "'answer' is the exact word or short term you removed, and it must "
                "NOT appear anywhere else in the question."
            )
        },
        {
            "type": "true-false",
            "shape": {"question": "<statement>", "answer": "<correct answer>", "options": None},
            "instructions": (
                "Write ONE original declarative STATEMENT, in your own words, that "
                "is clearly true or clearly false according to the content - not a "
                "sentence copied verbatim from it. This must be a statement, NOT a "
                "question: it must NOT end with a question mark, and it must NOT "
                "start with a word like 'Which', 'What', 'Does', 'Is', 'Are', or "
                "'Can'. For example, write 'X causes Y' instead of 'Does X cause Y?'. "
                "'answer' is exactly 'true' or 'false'."
            )
        },
        {
            "type": "short-answer",
            "shape": {"question": "<question>", "answer": "<short answer>", "options": None},
            "instructions": (
                "Ask a specific factual question that has exactly ONE unambiguous, "
                "short correct answer: a single word, number, or short phrase of at "
                "most 3-4 words (e.g. a name, term, or figure from the content). Do "
                "NOT ask a question whose correct answer would require a full "
                "sentence or is open to interpretation. 'answer' must be that short "
                "phrase, taken from the content - not a full sentence."
            )
        }
    ]

    asked = []
    topic_idx = 0
    for spec in qspecs:
        for _ in range(per_type):

            # assign a target topic to each question
            target_topic = topics[topic_idx % len(topics)] if topics else None
            topic_idx += 1

            # generate the question of the specific type, avoiding duplicates
            q = generate_ques(combined, spec, avoid=asked, topic=target_topic)
            if not q.get("failed"):
                q["topic"] = target_topic or "General"
                questions.append(q)
                asked.append(q["question"]) # keep track of asked questions to avoid duplicates

    return questions

@inst.timed("extract topics")
def extract_topics(text, n=8):
    """
    extract topics from source text
    """
    prompt = (
        "You are labelling the topics covered in a lecture.\n"
        f"List at most {n} distinct topics.\n"
        "Each topic must be 1-4 words and specific to the material.\n"
        'Respond with JSON only: {"topics": ["...", "..."]} \n\n'
        f"Lecture material:\n{text[:6000]}"
    )

    for _ in range(3):
        raw = ollama.generate(
            model="llama3.2", prompt=prompt, format="json",
            options={"temperature": 0.1, "seed": 42},
        )["response"]

        # extract topics from generated text
        try:
            topics = [str(t).strip() for t in json.loads(raw)["topics"] if str(t).strip()]
        except (json.JSONDecodeError, KeyError, TypeError):
            continue

        # dedupe the list of topics, keeping the order intact
        seen, out = set(), []
        for t in topics:
            if t.lower() not in seen:
                seen.add(t.lower())
                out.append(t)

        if out:
            return out[:n]

    return ["General"]

@inst.timed("generate score")
def generate_score(results):

    # consider the length of the question 
    # when considering the time spent on it
    times_per_word = [r["time"] / r["num_words"] for r in results]
    mean = statistics.mean(times_per_word)
    stdev = statistics.stdev(times_per_word) if len(times_per_word) > 1 else 0

    # assign incorrect/correct and confidence to each answer
    for r, tpw in zip(results, times_per_word):
        correct = str(r["user_answer"]).strip().lower() == str(r["correct_answer"]).strip().lower()
        r["correct"] = correct
        z = (tpw - mean) / stdev if stdev > 0 else 0
        r["confidence"] = 0.0 if not correct else 1 / (1 + math.exp(z))

    # keep track of confidence and score across topics
    topic_data = defaultdict(lambda: {"correct": 0, "total": 0, "confidences": []})
    for r in results:
        t = topic_data[r["topic"]]
        t["total"] += 1
        t["correct"] += int(r["correct"])
        t["confidences"].append(r["confidence"])

    # keep track of student's performance
    performance = []
    for topic, d in topic_data.items():
        performance.append({
            "topic": topic,
            "score": f"{d['correct']}/{d['total']}",
            "confidence": round(sum(d["confidences"]) / len(d["confidences"]), 2),
        })

    performance.sort(key=lambda x: x["confidence"])
    return json.dumps(performance, indent=2)

@inst.timed("generate feedback")
def generate_feedback(performance, combined):
    system = (
        "You write quiz feedback directly to the student, addressing them as "
        "'you'. You output ONLY the feedback itself - never a preamble, never "
        "a note to whoever asked for it, and never a sign-off."
    ) # ensure that feedback does not address prompter but only user

    prompt = f"""A student completed a quiz. Their performance per topic is provided
    (confidence is 0-1 where low means they were wrong or hesitant)
    Write encouraging, specific feedback directly to the student (use "you"):
    acknowledge strong topics and gently highlight weak ones. 2-3 short paragraphs.
    Stick strictly to the performance. Suggest improvements to weaknesses using the provided content.
    Do NOT include any preamble or meta-commentary such as "Here is your feedback"
    or "Based on the performance provided" - begin directly with the feedback.

    PERFORMANCE:
    {performance}

    CONTENT:
    {combined}
    """

    # generate the feedback according to performance
    feedback = ollama.chat(
        model="llama3.2",
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
    )["message"]["content"]

    return feedback

def _priority_bucket(confidence):
    """
    assign priority to topics based on confidence scores
    """
    if confidence < 0.4:
        return "high"
    if confidence < 0.7:
        return "medium"
    return "low"

def _allocate_days(topics, duration_days):
    """
    returns [(name, confidence, day_count), ...] with day_counts summing to
    duration_days, weighted towards the weakest topics.
    """
    n = len(topics)
    if n == 0 or duration_days <= 0:
        return []

    if duration_days <= n:
        return [(name, conf, 1) for name, conf in topics[:duration_days]]

    # deciding how many of the student's chosen study days go to each topic, 
    # weighted by how weak they are on it
    weights = [max(0.05, 1 - conf) for _, conf in topics]
    total_weight = sum(weights)
    quotas = [w / total_weight * duration_days for w in weights]
    days = [max(1, int(q)) for q in quotas]
    remainders = [q - int(q) for q in quotas]

    diff = duration_days - sum(days)
    if diff > 0:
        # hand out the leftover days to topics rounded down most
        order = sorted(range(n), key=lambda i: remainders[i], reverse=True)
        for i in order[:diff]:
            days[i] += 1
    else:
        # trim back starting from the topics closest to their rounded-down quota
        # since rounding up topics to 1-day minimum can exceed duration_days
        order = sorted(range(n), key=lambda i: remainders[i])
        i = 0
        while diff < 0:
            idx = order[i % n]
            if days[idx] > 1:
                days[idx] -= 1
                diff += 1
            i += 1

    return [(topics[i][0], topics[i][1], days[i]) for i in range(n)]

def _generate_topic_actions(performance, combined, retries=3):
    prompt = f"""A student completed a quiz. Their performance per topic is below
(confidence is 0-1, where low means they were wrong or hesitant).

PERFORMANCE:
{performance}

CONTENT:
{combined}

For EACH topic listed in PERFORMANCE, write 1 to 3 concrete, specific study
actions based ONLY on the CONTENT above (e.g. "re-read the section on X",
"practice deriving Y"). Order each topic's actions from most to least important.

Return ONLY JSON in this exact shape:
{{"actions": {{"<topic name exactly as given>": ["<action 1>", "<action 2>"]}}}}"""

    for retry in range(retries):
        response = ollama.chat(
            model="llama3.2",
            messages=[
                {"role": "system", "content": "You write specific study actions as JSON. You never write a paragraph."},
                {"role": "user", "content": prompt},
            ],
            format="json",
            options={"temperature": 0.3},
        ) # get content for the day for each topic in the learning plan

        try:
            data = json.loads(response["message"]["content"])
        except json.JSONDecodeError:
            print(f"Failed at attempt {retry}")
            continue

        actions = data.get("actions")
        if isinstance(actions, dict) and actions:

            # return the topic and its generated action
            return {
                str(topic): [str(a).strip() for a in items if str(a).strip()]
                for topic, items in actions.items()
                if isinstance(items, list) and items
            }

    return {}

@inst.timed("create plan")
def create_plan(performance, combined, duration_days=7):
    try:
        perf_list = json.loads(performance)
    except (json.JSONDecodeError, TypeError):
        perf_list = []

    # get each topic and its confidence
    topics = [(p["topic"], p["confidence"]) for p in perf_list if "topic" in p and "confidence" in p]
    if not topics:
        return []

    # get actions and allocate days
    actions_by_topic = _generate_topic_actions(performance, combined)
    allocation = _allocate_days(topics, duration_days)

    schedule = []
    day_num = 1
    for topic, confidence, day_count in allocation:

        # create the schedule organised according to the day and topic
        actions = actions_by_topic.get(topic) or ["Review your notes and the summary for this topic."]
        for i in range(day_count):
            schedule.append({
                "day": day_num,
                "topic": topic,
                "confidence": confidence,
                "priority": _priority_bucket(confidence),
                "action": actions[i % len(actions)],
            })
            day_num += 1

    return schedule

LLM_MODEL = "llama3.2"

def release_llm(model=LLM_MODEL, label="run", end_of_phase=True, unload=True):
    """
    free up GPU/RAM.
    unload=False keeps the model
    used between screens that reuse same model so it isn't reloaded each time
    """
    if unload:
        try:
            ollama.generate(model=model, prompt="", keep_alive=0)
        except Exception as exc:
            print(f"[release_llm] could not unload {model}: {exc}")
            return False

    if end_of_phase:
        inst.report()
        inst.save_run(label)
        inst.reset()

    return True