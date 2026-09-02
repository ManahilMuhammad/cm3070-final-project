import subprocess
import tempfile
from pathlib import Path
import cv2
import pdfplumber
from pptx import Presentation
from PIL import Image
import os
import time
import instrumentation as inst
import fitz
from io import BytesIO
from .utils import prepare_images

AUDIO_EXTENSIONS = ('.mp3', '.wav', '.m4a', '.aac', '.flac', '.wma', '.ogg')
VIDEO_EXTENSIONS = ('.mp4', '.mov', '.avi', '.mkv', '.flv', '.wmv')
DOCUMENT_EXTENSIONS = ('.pdf', '.pptx', '.ppt')
IMAGE_EXTENSIONS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp')

# --> AUDIO EXTRACTION
def extract_from_audio(audio_path):
    return str(audio_path)

# --> VIDEO EXTRACTION
def extract_from_video(video_path):
    """
    extract audio from video
    """
    video_path = str(video_path)
    
    # extract audio (copy file first)
    audio_path = tempfile.NamedTemporaryFile(suffix='.wav', delete=False).name

    result = subprocess.run([
        'ffmpeg',
        '-y', # overwrite output file                    
        '-i', video_path, # input file
        '-vn', # keep only audio                  
        '-acodec', 'pcm_s16le',
        audio_path # output file
    ], capture_output=True, text=True, timeout=300) # raise timeout if longer than 300 seconds taken

    if result.returncode != 0:
        raise RuntimeError(f'ffmpeg failed: {result.stderr}')
    
    if not os.path.exists(audio_path) or os.path.getsize(audio_path) == 0:
        raise RuntimeError(f'Audio extraction produced empty file: {audio_path}')
    
    return audio_path


# --> PDF EXTRACTION
def extract_from_pdf(pdf_path):
    """
    extract text and images from PDF
    """
    
    pdf_path = str(pdf_path)
    text_parts = []
    images = []
    
    doc = fitz.open(pdf_path)
    
    for page_num, page in enumerate(doc):
        # extract text
        page_text = page.get_text()
        if page_text.strip():
            text_parts.append(page_text)
        
        # extract images
        image_list = page.get_images(full=True)
        for img_index in image_list:
            xref = img_index[0]
            pix = fitz.Pixmap(doc, xref)
            
            # convert to PIL Image
            img_data = pix.tobytes('ppm')
            img = Image.open(BytesIO(img_data))
            img.load()  # force load into memory
            images.append(img)
            
            pix = None  # free memory
    
    doc.close()
    
    text = '\n\n'.join(text_parts)
    return text, images


# --> PPTX EXTRACTION
def extract_from_pptx(pptx_path):
    """
    extract text and images/figures from pptx
    """
    pptx_path = str(pptx_path)
    text_parts = []
    images = []
    
    prs = Presentation(pptx_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # extract text
            if hasattr(shape, 'text') and shape.text.strip():
                text_parts.append(shape.text)
            
            # extract images
            if shape.shape_type == 13:
                try:
                    image = shape.image

                    # convert python-pptx Image object to PIL Image
                    pil_image = Image.open(BytesIO(image.blob))
                    pil_image.load()

                    # make a fully independent PIL image
                    pil_image = pil_image.copy()
                    images.append(pil_image)

                except Exception as e:
                    print(f'Failed to extract image: {e}')
    
    text = '\n\n'.join(text_parts)
    return text, images


# --> UNIFIED ROUTER
def extract_from_file(file_obj):
    """
    route file to correct extraction function based on type
    returns:
        dict with keys: audio, text, images, notes
    """
    filename = file_obj.name.lower()
    
    result = {
        'audio': None,
        'text': '',
        'images': [],
        'notes': []
    }

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=Path(filename).suffix)
    tmp_path = tmp.name

    data = file_obj.getbuffer() if hasattr(file_obj, 'getbuffer') else file_obj.read()
    tmp.write(data)
    tmp.flush()
    tmp.close()
    
    try:
        if any(filename.endswith(ext) for ext in VIDEO_EXTENSIONS):
            # video
            audio = extract_from_video(tmp_path)
            result['audio'] = audio

        elif any(filename.endswith(ext) for ext in AUDIO_EXTENSIONS):
            # audio
            result['audio'] = extract_from_audio(tmp_path)
        
        elif filename.endswith('.pdf'):
            # PDF
            text, images = extract_from_pdf(tmp_path)
            result['text'] = text
            result['images'] = prepare_images(images)
        
        elif filename.endswith('.pptx'):
            # pptx
            text, images = extract_from_pptx(tmp_path)
            result['text'] = text
            result['images'] = prepare_images(images)
        
        elif any(filename.endswith(ext) for ext in IMAGE_EXTENSIONS):
            # single image (written notes)
            image = Image.open(tmp_path)
            result['notes'] = [image]
        
        else:
            raise ValueError(f'Unsupported file type: {filename}')
        
    finally:
        # give everything time to release the file
        time.sleep(0.1)
        
        # Now delete
        try:
            # clean up temp file
            # unless it is an audio file because that needs to be kept for transcription
            if os.path.exists(tmp_path) and not any(filename.endswith(ext) for ext in AUDIO_EXTENSIONS + VIDEO_EXTENSIONS):
                os.remove(tmp_path)
        except PermissionError:
            # if it still fails, let system cleanup on reboot
            pass
    
    return result